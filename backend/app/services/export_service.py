from docx import Document
from docx.shared import RGBColor, Pt as DocxPt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PptxRGBColor
from io import BytesIO
import markdown2
import re
from bs4 import BeautifulSoup
from typing import Dict, Any
import logging

# Configure logger with signature
logger = logging.getLogger(__name__)
logger.info("ExportService module loaded - Created by SAMBIT PRADHAN 22BCB0139")

class ExportService:

    # Define PPT Themes
    THEMES = {
        "professional": {
            "name": "professional",
            "display_name": "Professional & Corporate",
            "background_color": "#FFFFFF",
            "title_color": "#1F4788",
            "text_color": "#333333",
            "accent_color": "#2E75B6",
            "font_title": "Calibri",
            "font_body": "Calibri",
            "description": "Clean, minimal design with blue/gray color schemes"
        },
        "modern": {
            "name": "modern",
            "display_name": "Modern & Colorful",
            "background_color": "#F8F9FA",
            "title_color": "#6C5CE7",
            "text_color": "#2D3436",
            "accent_color": "#00B894",
            "font_title": "Arial",
            "font_body": "Arial",
            "description": "Vibrant colors and modern fonts"
        },
        "academic": {
            "name": "academic",
            "display_name": "Academic & Simple",
            "background_color": "#FEFEFE",
            "title_color": "#000000",
            "text_color": "#1A1A1A",
            "accent_color": "#4A4A4A",
            "font_title": "Times New Roman",
            "font_body": "Times New Roman",
            "description": "Traditional academic style with focus on content clarity"
        },
        "creative": {
            "name": "creative",
            "display_name": "Creative & Dynamic",
            "background_color": "#FFF5E6",
            "title_color": "#E74C3C",
            "text_color": "#2C3E50",
            "accent_color": "#F39C12",
            "font_title": "Arial",
            "font_body": "Arial",
            "description": "Bold colors and creative visual elements"
        }
    }

    @staticmethod
    def clean_html(raw_html):
        """Remove HTML tags but keep text content"""
        cleanr = re.compile('<.*?>')
        cleantext = re.sub(cleanr, '', raw_html)
        return cleantext

    @staticmethod
    def parse_html_with_formatting(html_content: str) -> list:
        """Parse HTML and extract formatted text segments with their styling"""
        if not html_content or '<' not in html_content:
            return [{'text': html_content, 'bold': False, 'italic': False}]

        soup = BeautifulSoup(html_content, 'html.parser')
        segments = []

        def extract_text_segments(element, parent_bold=False, parent_italic=False):
            """Recursively extract text with formatting"""
            if isinstance(element, str):
                text = str(element)
                if text and text.strip():  # Keep spaces but check if meaningful
                    segments.append({
                        'text': text,
                        'bold': parent_bold,
                        'italic': parent_italic
                    })
            else:
                # Check if this element adds formatting
                is_bold = parent_bold or element.name in ['strong', 'b']
                is_italic = parent_italic or element.name in ['em', 'i']

                # Process children
                for child in element.children:
                    extract_text_segments(child, is_bold, is_italic)

        extract_text_segments(soup)

        # Clean up: merge consecutive segments with same formatting
        if not segments:
            return [{'text': html_content, 'bold': False, 'italic': False}]

        cleaned = []
        for seg in segments:
            if cleaned and cleaned[-1]['bold'] == seg['bold'] and cleaned[-1]['italic'] == seg['italic']:
                # Merge with previous segment
                cleaned[-1]['text'] += seg['text']
            else:
                cleaned.append(seg)

        return cleaned

    @staticmethod
    def add_html_content_to_document(document, html_content):
        """Parse HTML and add formatted content to Word document with proper styling"""
        soup = BeautifulSoup(html_content, 'html.parser')

        def process_element(elem, list_level=0):
            """Process elements in order, handling nesting properly"""
            if elem.name == 'p':
                p = document.add_paragraph()
                ExportService._process_inline_elements(p, elem)
                # Set paragraph spacing
                p.paragraph_format.space_before = DocxPt(6)
                p.paragraph_format.space_after = DocxPt(6)
                p.paragraph_format.line_spacing = 1.15

            elif elem.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                level = min(int(elem.name[1]), 3)  # Word supports levels 1-9, we use 1-3
                heading = document.add_heading(elem.get_text(), level=level)
                heading.paragraph_format.space_before = DocxPt(12)
                heading.paragraph_format.space_after = DocxPt(6)

            elif elem.name in ['ul', 'ol']:
                for li in elem.find_all('li', recursive=False):
                    # Determine style based on list type and nesting
                    if elem.name == 'ul':
                        style = 'List Bullet' if list_level == 0 else 'List Bullet 2'
                    else:
                        style = 'List Number' if list_level == 0 else 'List Number 2'

                    p = document.add_paragraph(style=style)
                    ExportService._process_inline_elements(p, li)

                    # Set spacing for lists
                    p.paragraph_format.space_before = DocxPt(3)
                    p.paragraph_format.space_after = DocxPt(3)
                    p.paragraph_format.line_spacing = 1.15

                    # Process nested lists within this li
                    for nested_list in li.find_all(['ul', 'ol'], recursive=False):
                        process_element(nested_list, list_level + 1)

            elif elem.name == 'br':
                document.add_paragraph()

        # Process only top-level elements to avoid duplication
        for element in soup.find_all(recursive=False):
            process_element(element)

    @staticmethod
    def _process_inline_elements(paragraph, element):
        """Process inline HTML elements and apply formatting with proper font sizing"""
        for child in element.children:
            if isinstance(child, str):
                text = str(child)
                if text and text.strip():
                    run = paragraph.add_run(text)
                    run.font.size = DocxPt(11)  # Standard body text
                    run.font.name = 'Calibri'
            elif child.name == 'strong' or child.name == 'b':
                run = paragraph.add_run(child.get_text())
                run.bold = True
                run.font.size = DocxPt(11)
                run.font.name = 'Calibri'
            elif child.name == 'em' or child.name == 'i':
                run = paragraph.add_run(child.get_text())
                run.italic = True
                run.font.size = DocxPt(11)
                run.font.name = 'Calibri'
            elif child.name == 'u':
                run = paragraph.add_run(child.get_text())
                run.underline = True
                run.font.size = DocxPt(11)
                run.font.name = 'Calibri'
            elif child.name == 'span':
                run = paragraph.add_run(child.get_text())
                run.font.size = DocxPt(11)
                run.font.name = 'Calibri'
                style = child.get('style', '')
                if 'color:' in style:
                    color_match = re.search(r'color:\s*#([0-9a-fA-F]{6})', style)
                    if color_match:
                        hex_color = color_match.group(1)
                        r = int(hex_color[0:2], 16)
                        g = int(hex_color[2:4], 16)
                        b = int(hex_color[4:6], 16)
                        run.font.color.rgb = RGBColor(r, g, b)
            else:
                # Recursively process nested elements
                ExportService._process_inline_elements(paragraph, child)

    @staticmethod
    def generate_docx(project_data):
        logger.info(f"Starting DOCX generation for project: {project_data.get('title', 'Untitled')} - Created by SAMBIT PRADHAN 22BCB0139")
        document = Document()

        # Title Page
        document.add_heading(project_data.get('title', 'Untitled Project'), 0)
        document.add_paragraph(f"Generated by AI Doc Builder")
        document.add_paragraph(f"Created by SAMBIT PRADHAN 22BCB0139")
        document.add_page_break()

        # Sections
        for section in project_data.get('outline', []):
            document.add_heading(section.get('title', 'Untitled Section'), level=1)

            # REMOVED: Summary bullets are no longer included in DOCX export
            # They remain visible in the UI but are excluded from downloads

            # Add main content with HTML formatting
            content = section.get('content', '')
            if content:
                # Check if content is HTML (contains tags)
                if '<' in content and '>' in content:
                    # Parse and format HTML content
                    ExportService.add_html_content_to_document(document, content)
                else:
                    # Fallback for plain text or markdown
                    html_content = markdown2.markdown(content)
                    text_content = ExportService.clean_html(html_content)
                    document.add_paragraph(text_content)

            document.add_paragraph() # Spacing

        # Footer
        section = document.sections[0]
        footer = section.footer
        paragraph = footer.paragraphs[0]
        paragraph.text = "AI Doc Builder Export"

        f = BytesIO()
        document.save(f)
        f.seek(0)
        return f

    @staticmethod
    def hex_to_rgb(hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    @staticmethod
    def get_theme(theme_name: str = "professional") -> Dict[str, Any]:
        """Get theme configuration by name"""
        return ExportService.THEMES.get(theme_name, ExportService.THEMES["professional"])

    @staticmethod
    def calculate_font_size(text_length: int, base_size: int = 18) -> int:
        """Calculate appropriate font size based on text length"""
        if text_length < 200:
            return base_size
        elif text_length < 400:
            return max(base_size - 2, 14)
        elif text_length < 600:
            return max(base_size - 4, 12)
        else:
            return max(base_size - 6, 10)

    @staticmethod
    def has_bullet_markers(text: str) -> bool:
        """Check if text contains bullet point markers"""
        clean_text = ExportService.clean_html(text) if '<' in text else text
        lines = clean_text.strip().split('\n')

        bullet_patterns = [
            r'^\s*[-•*]\s+',  # Dash, bullet, asterisk
            r'^\s*\d+[\.)]\s+',  # Numbered lists (1. or 1))
        ]

        bullet_count = 0
        for line in lines:
            line = line.strip()
            if line:
                for pattern in bullet_patterns:
                    if re.match(pattern, line):
                        bullet_count += 1
                        break

        # If more than 2 lines have bullet markers, treat as bullet list
        return bullet_count >= 2

    @staticmethod
    def extract_bullets_from_text(text: str) -> list:
        """Extract bullet points from text that has bullet markers"""
        clean_text = ExportService.clean_html(text) if '<' in text else text
        lines = clean_text.strip().split('\n')

        bullets = []
        for line in lines:
            line = line.strip()
            if line:
                # Remove bullet markers
                line = re.sub(r'^\s*[-•*]\s+', '', line)
                line = re.sub(r'^\s*\d+[\.)]\s+', '', line)
                if line:
                    bullets.append(line)

        return bullets[:5]  # Limit to 5 bullets

    @staticmethod
    def parse_html_to_structured_content(html_content: str, max_chars: int = 1000) -> list:
        """Parse HTML content into structured paragraphs/bullets with formatting preserved"""
        if not html_content or '<' not in html_content:
            return [{'type': 'paragraph', 'content': html_content, 'text': html_content, 'level': 0}]

        soup = BeautifulSoup(html_content, 'html.parser')
        items = []
        total_chars = 0

        def process_element(elem, parent_level=0):
            nonlocal total_chars

            if total_chars >= max_chars:
                return

            if elem.name == 'p':
                text = elem.get_text().strip()
                if text and total_chars + len(text) <= max_chars:
                    items.append({
                        'type': 'paragraph',
                        'content': str(elem),
                        'text': text,
                        'level': 0
                    })
                    total_chars += len(text)

            elif elem.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                text = elem.get_text().strip()
                if text and total_chars + len(text) <= max_chars:
                    items.append({
                        'type': 'heading',
                        'content': str(elem),
                        'text': text,
                        'level': 0
                    })
                    total_chars += len(text)

            elif elem.name in ['ul', 'ol']:
                # Process all direct li children
                for li in elem.find_all('li', recursive=False):
                    if total_chars >= max_chars:
                        break

                    text = li.get_text().strip()
                    if text and total_chars + len(text) <= max_chars:
                        items.append({
                            'type': 'bullet',
                            'content': str(li),
                            'text': text,
                            'level': parent_level
                        })
                        total_chars += len(text)

                        # Check for nested lists
                        for nested in li.find_all(['ul', 'ol'], recursive=False):
                            process_element(nested, parent_level + 1)

        # Process top-level elements only
        for element in soup.find_all(recursive=False):
            process_element(element)

        return items if items else [{'type': 'paragraph', 'content': html_content, 'text': ExportService.clean_html(html_content), 'level': 0}]

    @staticmethod
    def format_as_paragraphs(text: str, max_chars: int = 800) -> list:
        """Fallback: Format text as paragraph chunks (used when no HTML)"""
        clean_text = ExportService.clean_html(text) if '<' in text else text

        # Split by double newlines or single newlines
        paragraphs = [p.strip() for p in clean_text.split('\n') if p.strip()]

        if not paragraphs:
            paragraphs = [clean_text.strip()]

        # Smart truncation
        result = []
        total_chars = 0

        for para in paragraphs:
            if total_chars + len(para) <= max_chars:
                result.append(para)
                total_chars += len(para)
            elif total_chars < max_chars * 0.7:
                remaining = max_chars - total_chars
                if remaining > 50:
                    result.append(para[:remaining] + "...")
                break
            else:
                break

        return result if result else paragraphs[:5]

    @staticmethod
    def generate_pptx(project_data, theme_name: str = "professional"):
        """Generate PowerPoint with proper formatting and theme support"""
        logger.info(f"Starting PPTX generation for project: {project_data.get('title', 'Untitled')} with theme {theme_name} - Created by SAMBIT PRADHAN 22BCB0139")
        prs = Presentation()
        theme = ExportService.get_theme(theme_name)

        # Set slide dimensions to standard size (10" x 7.5")
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # ========== TITLE SLIDE ==========
        blank_layout = prs.slide_layouts[6]  # Blank layout for custom design
        title_slide = prs.slides.add_slide(blank_layout)

        # Background
        background = title_slide.background
        fill = background.fill
        fill.solid()
        bg_rgb = ExportService.hex_to_rgb(theme["background_color"])
        fill.fore_color.rgb = PptxRGBColor(*bg_rgb)

        # Get color components
        title_rgb = ExportService.hex_to_rgb(theme["title_color"])
        text_rgb = ExportService.hex_to_rgb(theme["text_color"])
        accent_rgb = ExportService.hex_to_rgb(theme["accent_color"])

        # Top accent bar
        top_bar = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.3)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = PptxRGBColor(*accent_rgb)
        top_bar.line.fill.background()

        # Bottom accent bar
        bottom_bar = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(7.2), Inches(10), Inches(0.3)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = PptxRGBColor(*accent_rgb)
        bottom_bar.line.fill.background()

        # Title box
        title_left = Inches(1)
        title_top = Inches(2.5)
        title_width = Inches(8)
        title_height = Inches(1.5)

        title_box = title_slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_para = title_frame.paragraphs[0]
        title_para.text = project_data.get('title', 'Untitled Project')
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.name = theme["font_title"]
        title_para.font.color.rgb = PptxRGBColor(*title_rgb)

        # Decorative line under title
        line_shape = title_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3), Inches(4.1), Inches(4), Inches(0.05)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = PptxRGBColor(*accent_rgb)
        line_shape.line.fill.background()

        # Subtitle
        subtitle_top = Inches(4.4)
        subtitle_box = title_slide.shapes.add_textbox(title_left, subtitle_top, title_width, Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True

        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = "Generated by AI Doc Builder"
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.name = theme["font_body"]
        subtitle_para.font.color.rgb = PptxRGBColor(*text_rgb)

        # Add signature
        signature_para = subtitle_frame.add_paragraph()
        signature_para.text = "Created by SAMBIT PRADHAN 22BCB0139"
        signature_para.alignment = PP_ALIGN.CENTER
        signature_para.font.size = Pt(14)
        signature_para.font.name = theme["font_body"]
        signature_para.font.color.rgb = PptxRGBColor(*text_rgb)

        # ========== CONTENT SLIDES ==========
        for idx, section in enumerate(project_data.get('outline', []), 1):
            slide = prs.slides.add_slide(blank_layout)

            # Background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = PptxRGBColor(*bg_rgb)

            # LEFT ACCENT STRIPE (vertical colored bar)
            left_stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(0.15), Inches(7.5)
            )
            left_stripe.fill.solid()
            left_stripe.fill.fore_color.rgb = PptxRGBColor(*accent_rgb)
            left_stripe.line.fill.background()

            # TITLE BAR BACKGROUND (colored rectangle behind title)
            title_bar_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.15), Inches(0.3), Inches(9.85), Inches(0.9)
            )
            title_bar_bg.fill.solid()
            title_bar_bg.fill.fore_color.rgb = PptxRGBColor(*accent_rgb)
            title_bar_bg.line.fill.background()

            # TITLE TEXT (white text on colored background)
            title_left = Inches(0.4)
            title_top = Inches(0.4)
            title_width = Inches(9.2)
            title_height = Inches(0.7)

            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.margin_bottom = Pt(5)
            title_frame.margin_top = Pt(5)
            title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            title_para = title_frame.paragraphs[0]
            title_para.text = section.get('title', 'Untitled')
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.font.name = theme["font_title"]
            # White text on colored background for contrast
            title_para.font.color.rgb = PptxRGBColor(255, 255, 255)

            # DECORATIVE DIVIDER LINE (under title bar)
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.4), Inches(1.35), Inches(9.2), Inches(0.03)
            )
            divider.fill.solid()
            # Use title color with some transparency effect
            divider.fill.fore_color.rgb = PptxRGBColor(*title_rgb)
            divider.line.fill.background()

            # Content area
            content_left = Inches(0.7)
            content_top = Inches(1.5)
            content_width = Inches(8.6)
            content_height = Inches(5.3)

            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_frame.margin_left = Pt(10)
            content_frame.margin_right = Pt(10)
            content_frame.margin_top = Pt(10)
            content_frame.margin_bottom = Pt(10)

            # PARSE HTML CONTENT WITH FORMATTING PRESERVED
            content = section.get('content', '')

            if content:
                # Parse HTML into structured items with formatting
                structured_items = ExportService.parse_html_to_structured_content(content)

                # DYNAMIC FONT SIZING based on content length
                total_length = sum(len(item.get('text', '')) for item in structured_items)
                num_items = len(structured_items)

                if total_length > 800 or num_items > 7:
                    font_size = 14
                elif total_length > 600 or num_items > 5:
                    font_size = 16
                elif total_length > 400 or num_items > 4:
                    font_size = 18
                else:
                    font_size = 20

                # Render each structured item
                for i, item in enumerate(structured_items):
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()

                    item_type = item.get('type', 'paragraph')
                    item_level = item.get('level', 0)
                    html_content = item.get('content', '')

                    # Set bullet level for indentation
                    p.level = item_level

                    # Parse HTML formatting and add runs
                    if html_content and '<' in html_content:
                        segments = ExportService.parse_html_with_formatting(html_content)

                        for seg in segments:
                            run = p.add_run()
                            run.text = seg['text']
                            run.font.size = Pt(font_size)
                            run.font.name = theme["font_body"]
                            run.font.bold = seg.get('bold', False)
                            run.font.italic = seg.get('italic', False)

                            # Color: bold text uses title color for emphasis
                            if seg.get('bold'):
                                run.font.color.rgb = PptxRGBColor(*title_rgb)
                            else:
                                run.font.color.rgb = PptxRGBColor(*text_rgb)
                    else:
                        # Plain text fallback
                        p.text = item.get('text', '')
                        p.font.size = Pt(font_size)
                        p.font.name = theme["font_body"]
                        p.font.color.rgb = PptxRGBColor(*text_rgb)

                    # Spacing based on type
                    if item_type == 'bullet':
                        p.space_before = Pt(4)
                        p.space_after = Pt(4)
                        p.line_spacing = 1.2
                    elif item_type == 'heading':
                        p.space_before = Pt(12)
                        p.space_after = Pt(6)
                        p.line_spacing = 1.1
                        # Make heading bold if not already
                        if p.runs:
                            for run in p.runs:
                                run.font.bold = True
                                run.font.size = Pt(font_size + 2)
                    else:  # paragraph
                        p.space_before = Pt(6)
                        p.space_after = Pt(8)
                        p.line_spacing = 1.3

            # Footer with slide number
            footer_top = Inches(7.1)
            footer_box = slide.shapes.add_textbox(Inches(8.5), footer_top, Inches(1), Inches(0.3))
            footer_frame = footer_box.text_frame
            footer_para = footer_frame.paragraphs[0]
            footer_para.text = f"{idx}"
            footer_para.alignment = PP_ALIGN.RIGHT
            footer_para.font.size = Pt(12)
            footer_para.font.name = theme["font_body"]
            accent_rgb = ExportService.hex_to_rgb(theme["accent_color"])
            footer_para.font.color.rgb = PptxRGBColor(*accent_rgb)

            # Add full content to notes
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = ExportService.clean_html(content) if content else ''

        f = BytesIO()
        prs.save(f)
        f.seek(0)
        return f
